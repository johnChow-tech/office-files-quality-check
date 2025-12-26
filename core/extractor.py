# ==============================================================================
# Office ファイルからコンテンツを抽出するコアビジネスロジック (進捗コールバック機能を追加済み)
# ==============================================================================
import os
import csv
import sys
from datetime import datetime

# 異なるファイルタイプを処理するためのライブラリをインポート (インストールを確認: python-docx openpyxl python-pptx)
try:
    from docx import Document
    from openpyxl import load_workbook
    from pptx import Presentation
except ImportError:
    # このエラーメッセージは現在、ログウィンドウにもリダイレクトされます
    print("エラー: 必要な Office ファイル処理ライブラリが不足しています。pip install python-docx openpyxl python-pptx を実行してください", file=sys.stderr)


class DocExtractor:
    """

    コアエクストラクタクラス。指定されたパスの Office ファイルからプレーンテキストとハイパーリンクを抽出します。
    """
    # サポートされているファイルタイプを定義 (新たに .xlsm を追加)
    TARGET_EXTENSIONS = ('.docx', '.xlsx', '.xlsm', '.pptx')
    TEXT_EXTENSIONS = ('.md')
    URLS_EXTENSIONS = ('.dat')

    # 出力サブフォルダ名を定義
    formatted_time = datetime.now().strftime("%Y%m%d%H%M%S")
    TEXT_FOLDER = f"PlainText_{formatted_time}"
    URLS_FOLDER = f"HyperLinks_{formatted_time}"

    def __init__(self):
        print("DocExtractor の初期化が完了しました。")

    # <<< run_extractionの変更点：progress_callbackパラメータの追加 >>>
    def run_extraction(self, source_path: str, output_path: str, file_map: dict, progress_callback=None):
        """
        GUI から提供されたファイルリストとパスに基づいて、抽出と保存の全プロセスを実行します。
        progress_callback: GUI の進捗を更新するためのコールバック関数です。
        """
        if not source_path or not output_path or not file_map:
            print("エラー：ソースパス、出力パス、またはファイルリストは空にできません。", file=sys.stderr)
            return False

        print(f"抽出を開始します。ソースパス: {source_path}, 出力パス: {output_path}")

        # 1. 出力ディレクトリ構造が存在することを確認
        text_output_dir = os.path.join(output_path, self.TEXT_FOLDER)
        urls_output_dir = os.path.join(output_path, self.URLS_FOLDER)

        try:
            os.makedirs(text_output_dir, exist_ok=True)
            os.makedirs(urls_output_dir, exist_ok=True)
        except Exception as e:
            # ここでのエラーは呼び出し元のスレッドで処理され、GUIに返されます
            raise IOError(f"出力ディレクトリを作成できません: {e}")

        # 2. ファイルリストを走査 (file_map には処理対象ファイルの絶対パスが格納されています)
        total_files = len(file_map)
        processed_count = 0

        # 進捗バーが開始時に少なくとも一度呼び出されることを確認
        if progress_callback:
            progress_callback(processed_count, total_files)

        for index, file_abs_path in file_map.items():
            file_name = os.path.basename(file_abs_path)
            print(f"処理中 ({processed_count+1}/{total_files}): {file_name}")

            # _process_single_file を呼び出します。これは内部でファイル処理とエラー分離を行います
            self._process_single_file(
                file_abs_path, text_output_dir, urls_output_dir)

            # <<< 進捗の更新 >>>
            processed_count += 1
            if progress_callback:
                progress_callback(processed_count, total_files)

        print("抽出プロセスが終了しました。")
        return True

    def _process_single_file(self, file_abs_path, text_output_dir, urls_output_dir):
        """
        単一ファイルを処理し、対応する抽出メソッドを呼び出し、結果を保存します。
        """

        # 1. ファイル名と拡張子の解析
        original_file_name = os.path.basename(file_abs_path)
        base_name, ext = os.path.splitext(original_file_name)
        ext = ext.lower()   # 统一小写，方便判断
        # ext = ext.lower()   # 判定しやすいように小文字に統一

        text_content = ""
        urls_list = []

        try:
            # 2. コア抽出ロジック: ファイルのオープンエラーとタイプエラーを処理
            if ext == '.docx':
                print(f"  [DOCX] 抽出中です...")
                text_content, urls_list = self._extract_from_docx(
                    file_abs_path)

            elif ext == '.xlsx' or ext == '.xlsm':
                print(f"  [XLSX/XLSM] 抽出中です...")
                text_content, urls_list = self._extract_from_xlsx(
                    file_abs_path)

            elif ext == '.pptx':
                print(f"  [PPTX] 抽出中です...")
                text_content, urls_list = self._extract_from_pptx(
                    file_abs_path)

            else:
                print(f"  [スキップ] サポートされていないファイルタイプ: {ext}")
                return

        except Exception as e:
            # ファイルの読み取り時やオープン時の致命的なエラーをキャッチ (例: "File is not a zip file")
            print(
                f"    [エラー] ファイル {original_file_name} の処理中に致命的な読み取り/オープン例外が発生しました: {e}", file=sys.stderr)
            return

        # 3. プレーンテキストの保存を試行 (保存エラーは個別に処理され、ハイパーリンクの保存には影響しません)
        try:
            if text_content:
                self._save_plain_text(original_file_name,
                                      text_content, text_output_dir)
        except Exception as e:
            print(
                f"  [エラー] プレーンテキスト {original_file_name} を保存できません: {e}", file=sys.stderr)

        # 4. ハイパーリンクの保存を試行 (保存エラーは個別に処理され、プレーンテキストの保存には影響しません)
        try:
            if urls_list:
                self._save_urls_to_csv(
                    original_file_name, urls_list, urls_output_dir)
        except Exception as e:
            print(
                f"  [エラー] ハイパーリンク CSV {original_file_name} を保存できません: {e}", file=sys.stderr)

    # --------------------------------------------------------------------------
    #             ファイル抽出ロジック (A 部)
    # --------------------------------------------------------------------------

    def _extract_from_docx(self, file_abs_path):
        # DOCX ファイルからプレーンテキストとハイパーリンクを抽出します。
        try:
            document = Document(file_abs_path)
        except Exception as e:
            raise IOError(f"DOCX ファイルを開くまたは読み取ることができません: {e}")

        full_text = []
        urls_list = []

        # 段落を走査
        for paragraph in document.paragraphs:
            full_text.append(paragraph.text)

            # ハイパーリンクを抽出 (安全チェック済み)
            for run in paragraph.runs:
                try:
                    if hasattr(run, 'hyperlink') and run.hyperlink and run.hyperlink.address:
                        url = run.hyperlink.address
                        link_text = run.text.strip()

                        if url and link_text:
                            urls_list.append((link_text, url))

                except KeyError as ke:
                    print(f"  [DOCX警告] 無効な、または内部的なハイパーリンクを検出しました: {ke}")
                except AttributeError:
                    pass
                except Exception as e:
                    print(f"  [DOCX警告] ハイパーリンクの抽出中に不明なエラーが発生しました: {e}")

        # 表を走査
        for table in document.tables:
            for row in table.rows:
                for cell in row.cells:
                    for paragraph in cell.paragraphs:
                        full_text.append(paragraph.text)

        text_content = '\n'.join(full_text)
        unique_urls = list(set(urls_list))

        return text_content, unique_urls

    def _extract_from_xlsx(self, file_abs_path):
        # XLSX/XLSM ファイルからプレーンテキストとハイパーリンクを抽出します。
        try:
            workbook = load_workbook(
                file_abs_path, read_only=True, data_only=True)
        except Exception as e:
            raise IOError(f"XLSX ファイルを開くまたは読み取ることができません: {e}")

        full_text = []
        urls_list = []

        for sheet in workbook.worksheets:
            full_text.append(f"\n--- Sheet: {sheet.title} ---\n")

            for row in sheet.iter_rows():
                row_content = []
                for cell in row:
                    cell_value = str(
                        cell.value) if cell.value is not None else ""
                    row_content.append(cell_value)

                    # ハイパーリンクを抽出 (安全チェック済み)
                    try:
                        if hasattr(cell, 'hyperlink') and cell.hyperlink and cell.hyperlink.target:
                            url = cell.hyperlink.target
                            link_text = cell_value

                            if url and link_text:
                                urls_list.append((link_text, url))
                    except AttributeError:
                        pass
                    except Exception as e:
                        print(f"  [XLSX警告] ハイパーリンクの抽出中に不明なエラーが発生しました: {e}")

                if any(row_content):
                    full_text.append('\t'.join(row_content))

        text_content = '\n'.join(full_text)
        unique_urls = list(set(urls_list))

        return text_content, unique_urls

    def _extract_from_pptx(self, file_abs_path):
        # PPTX ファイルからプレーンテキストとハイパーリンクを抽出します。
        try:
            presentation = Presentation(file_abs_path)
        except Exception as e:
            raise IOError(f"PPTX ファイルを開くまたは読み取ることができません: {e}")

        full_text = []
        urls_list = []

        for slide_index, slide in enumerate(presentation.slides):
            full_text.append(f"\n--- Slide {slide_index + 1} ---\n")

            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue

                for paragraph in shape.text_frame.paragraphs:
                    para_text = ""

                    for run in paragraph.runs:
                        para_text += run.text

                        # ハイパーリンクを抽出 (docx と同様、安全チェック済み)
                        try:
                            if hasattr(run, 'hyperlink') and run.hyperlink and run.hyperlink.address:
                                url = run.hyperlink.address
                                link_text = run.text.strip()

                                if url and link_text:
                                    urls_list.append((link_text, url))
                        except AttributeError:
                            pass
                        except Exception as e:
                            print(f"  [PPTX警告] ハイパーリンクの抽出中に不明なエラーが発生しました: {e}")

                    if para_text.strip():
                        full_text.append(para_text)

        text_content = '\n'.join(full_text)
        unique_urls = list(set(urls_list))

        return text_content, unique_urls

    # --------------------------------------------------------------------------
    #             ファイル保存ヘルパーメソッド (命名を修正)
    # --------------------------------------------------------------------------

    def _save_plain_text(self, original_file_name, content, output_dir):
        """
        プレーンテキストのコンテンツを Markdown (.md) ファイルに保存します。
        ターゲット形式: PlainText_[ファイル名]_[拡張子].md
        """
        # 1. 元のファイル名を解析: '報告書.docx' -> ('報告書', '.docx')
        base_name, ext = os.path.splitext(original_file_name)
        original_ext_clean = ext.strip('.')  # 'docx'
        # original_ext_clean = ext.strip('.')  # 'docx'
        # 2. 新しいファイル名を構築: PlainText_報告書_docx.md
        output_file_name = f"PlainText_{base_name}_{original_ext_clean}.md"
        output_file = os.path.join(output_dir, output_file_name)

        try:
            with open(output_file, 'w', encoding='utf-8') as f:
                f.write(content)
            print(f"  -> テキストは {os.path.basename(output_file)} に保存されました")
        except Exception as e:
            raise e  # 抛出错误，交给 _process_single_file 捕获
            # raise e  # エラーを発生させ、_process_single_file にキャッチさせる

    def _save_urls_to_csv(self, original_file_name, urls_list, output_dir):
        """
        ハイパーリンクのリストを CSV ファイルに保存します。
        ターゲット形式: Urls_[ファイル名]_[拡張子].csv
        """
        # 1. 元のファイル名を解析: '報告書.docx' -> ('報告書', '.docx')
        base_name, ext = os.path.splitext(original_file_name)
        original_ext_clean = ext.strip('.')  # 'docx'
        # original_ext_clean = ext.strip('.')  # 'docx'

        # 2. 新しいファイル名を構築: Urls_報告書_docx.csv
        output_file_name = f"Urls_{base_name}_{original_ext_clean}.dat"
        output_file = os.path.join(output_dir, output_file_name)

        try:
            with open(output_file, 'w', newline='', encoding='utf-8') as csvfile:
                writer = csv.writer(csvfile)
                writer.writerow(["Source File", "Link Text", "URL"])  # ヘッダー行

                for link_text, url in urls_list:
                    # CSV 内部の Source File フィールドには、元のファイル名（拡張子付き）を保持
                    writer.writerow([original_file_name, link_text, url])
            print(f"  -> ハイパーリンクは {os.path.basename(output_file)} に保存されました")
        except Exception as e:
            raise e  # エラーを発生させ、_process_single_file にキャッチさせる
